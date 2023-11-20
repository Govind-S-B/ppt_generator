from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE

import re
import io


def sanitize_string(input_str):
    # Remove non-alphanumeric, underscores, hyphens, and periods
    sanitized = re.sub(r"[^A-Za-z0-9_.-]", "", input_str)

    # Replace consecutive periods with a single period
    sanitized = re.sub(r"\.{2,}", ".", sanitized)

    # Ensure the string starts and ends with an alphanumeric character
    sanitized = re.sub(r"^[^A-Za-z0-9]+", "", sanitized)
    sanitized = re.sub(r"[^A-Za-z0-9]+$", "", sanitized)

    # Truncate or pad string to meet the 3-63 character length requirement
    sanitized = sanitized[:63] if len(
        sanitized) > 63 else sanitized.ljust(3, "_")

    return sanitized


def ppt_gen(slide_data):
    ppt = Presentation()

    # Setting Background
    slide_master = ppt.slide_master
    slide_master.background.fill.solid()
    slide_master.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Title Screen
    curr_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    curr_slide.shapes.title.text = slide_data[0][0]
    curr_slide.shapes.title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    curr_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255)
    curr_slide.shapes.placeholders[1].text = slide_data[0][1]
    curr_slide.shapes.placeholders[1].text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255)

    # Overview
    curr_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    curr_slide.shapes.title.text = "Overview"
    curr_slide.shapes.title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    curr_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255)
    for content in slide_data[1]:
        tframe = curr_slide.shapes.placeholders[1].text_frame
        tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        para = tframe.add_paragraph()
        para.text = content
        para.level = 1
        para.font.color.rgb = RGBColor(
            255, 255, 255)

    # Content Slides
    for curr_slide_data in slide_data[2:]:
        curr_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        curr_slide.shapes.title.text = curr_slide_data[0]
        curr_slide.shapes.title.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        curr_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(
            255, 255, 255)
        for content in curr_slide_data[1:]:
            tframe = curr_slide.shapes.placeholders[1].text_frame
            tframe.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            para = tframe.add_paragraph()
            para.text = content
            para.level = 1
            para.font.color.rgb = RGBColor(
                255, 255, 255)

    # Thank You Screen
    curr_slide = ppt.slides.add_slide(ppt.slide_layouts[2])
    curr_slide.shapes.placeholders[1].text = "Thank You"

    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(
        255, 255, 255)
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].font.size = Pt(
        96)
    curr_slide.shapes.placeholders[1].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # f"{sanitize_string(slide_data[0][0])}.pptx"
    ppt_stream = io.BytesIO()
    ppt.save(ppt_stream)
    ppt_stream.seek(0)

    return ppt_stream
